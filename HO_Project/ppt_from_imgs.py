from pptx import Presentation
import os
import glob




os.chdir(r'C:\Users\ksraj\OneDrive - Indian Oil Corporation Limited\HO_10ROS_Each_25_Photos')

x=os.listdir()

for i in range(len(x)):
    print(x[i])
    os.chdir("C:\\Users\\ksraj\\OneDrive - Indian Oil Corporation Limited\\HO_10ROS_Each_25_Photos\\"+x[i])
    y=os.listdir()
    # print(len(y))

#
# os.chdir("C:\\Users\\ksraj\\OneDrive - Indian Oil Corporation Limited\\HO_10ROS_Each_25_Photos\\Kozhikode I\\")
# y = os.listdir()

    for j in range(len(y)):

        parent_dir="C:\\Users\\ksraj\\OneDrive - Indian Oil Corporation Limited\\HO_10ROS_Each_25_Photos\\"+x[i]+"\\"+y[j]
        directory="new\\"
        path=os.path.join(parent_dir,directory)
        print(path)
        # os.mkdir(path)
        os.chdir(path)
        z=os.listdir()


# os.chdir(r'C:\Users\ksraj\OneDrive - Indian Oil Corporation Limited\HO_10ROS_Each_25_Photos\Kanhangad\ALSAFAR FUELS\new')
#         x=os.listdir()
#         print(len(x))
        imgs=glob.glob("*.png")
        prs=Presentation('newPPT.pptx')
        for k in range(len(imgs)):
            layout=prs.slide_layouts[3]
            slide=prs.slides.add_slide(layout)
            title=slide.shapes.title.text="RO photos"
            # for shape in slide.placeholders:
            #     print('%d %s' % (shape.placeholder_format.idx, shape.name))

            slide.placeholders[1].text="RO Before Photo"
            slide.placeholders[15].text="RO After Photo"
            img=imgs[k]
            slide.placeholders[13].insert_picture(img)
            #
            prs.save('newPPT.pptx')
        # os.startfile("newPPT.pptx")

        # for i in range(len(x)):
        #
        #     folder_name = 'C:\\Users\\ksraj\\OneDrive - Indian Oil Corporation Limited\\HO_10ROS_Each_25_Photos\\Tirur\\' +x[i]
        #     os.chdir(folder_name)
        #     pr1 = Presentation()
        #     pr1.save("newPPT.pptx")
        #
        # # os.chdir(r'C:\Users\k sraj\OneDrive - Indian Oil Corporation Limited\HO_10ROS_Each_25_Photos\Kanhangad\Chandragiri fuels')
        #
        #

        # pr1.save("newPPT.pptx")
